#!/usr/bin/env python3
"""Build iAssetsPro WO Workflow Presentation with real VPS screenshots."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Config ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SCREENSHOTS = "/home/z/my-project/vps-screenshots"
OUTPUT = "/home/z/my-project/public/iAssetsPro-WO-Workflow-Presentation.pptx"

# Colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Deep navy
BG_CARD = RGBColor(0x1A, 0x23, 0x3B)        # Card background
ACCENT = RGBColor(0x10, 0xB9, 0x81)         # Emerald green
ACCENT2 = RGBColor(0x06, 0xB6, 0xD4)        # Cyan
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
DIM = RGBColor(0x94, 0xA3, 0xB8)
STEP_NUM = RGBColor(0x34, 0xD3, 0x99)       # Light emerald for step numbers

# ── Slide definitions ──────────────────────────────────────────────
# (file, title, subtitle, step_num, description_bullets)
slides_data = [
    ("00-login.png", "Login & Authentication",
     "Secure role-based access with demo accounts",
     "01",
     ["Users log in with username/password credentials",
      "Role-based access control (Operator, Supervisor, Admin, etc.)",
      "Each role has specific permissions for the WO workflow",
      "Demo accounts available for quick access to different roles"]),

    ("02-admin-dashboard.png", "Dashboard Overview",
     "Central command center with real-time KPIs",
     "02",
     ["Displays active Work Orders, pending requests, and overdue items",
      "Quick actions: New Request, View WOs, Approvals, PM Schedules",
      "Recent maintenance requests and work orders at a glance",
      "Cross-module overview: PM Overdue, Assets at Risk"]),

    ("03-maintenance-requests.png", "Maintenance Requests List",
     "View, filter, and manage all maintenance requests",
     "03",
     ["Comprehensive table with request #, title, priority, status",
      "Filter by status (Pending, Approved, Converted) and priority",
      "Shows requester name, associated asset, and date",
      "Click any row to open the full request detail in a side panel"]),

    ("04-create-request-dialog.png", "Create Maintenance Request",
     "Operator initiates a new maintenance request",
     "04",
     ["Title and detailed description of the issue",
      "Machine/asset selection with search capability",
      "Priority level (Low, Medium, High, Urgent)",
      "Machine downtime status and request type selection"]),

    ("05-work-orders.png", "Work Orders Management",
     "Track and manage all work orders across the organization",
     "05",
     ["Work orders created from approved maintenance requests",
      "Type classification: Breakdown, Corrective, Preventive",
      "Status tracking: Assigned, In Progress, Completed, Closed",
      "Assigned technician and creation date visibility"]),

    ("06-wo-detail.png", "Work Order Execution & Detail",
     "Detailed view with actions: Start Work, Complete WO",
     "06",
     ["Full WO detail in a slide-over panel",
      "Action buttons: Start Work, Complete WO, Request Tools/Materials",
      "Task checklist and time logging capabilities",
      "Team member management and assistance requests"]),

    ("07-tool-requests.png", "Repairs — Tool Requests",
     "Request specialized tools needed for WO execution",
     "07",
     ["Technicians request tools required for specific work orders",
      "Linked to WO number for traceability",
      "Request status tracking: Pending, Approved, Issued",
      "Integration with tool room inventory management"]),

    ("08-material-requests.png", "Repairs — Material Requests",
     "Request spare parts and materials for repairs",
     "08",
     ["Spare parts and consumables requested against WOs",
      "Material description, quantity, and unit specification",
      "Approval workflow before issuance from stores",
      "Full traceability from request to consumption"]),

    ("09-tool-transfers.png", "Repairs — Tool Transfers",
     "Transfer tools between plants or departments",
     "09",
     ["Inter-plant and inter-department tool transfers",
      "Transfer tracking with sender and receiver details",
      "Ensures tools are available at the right location",
      "Complete audit trail for tool movement"]),

    ("10-completion.png", "Repairs — WO Completion",
     "Complete work orders with final documentation",
     "10",
     ["Mark work orders as completed after finishing repairs",
      "Document actual work performed and time spent",
      "Record parts used and tools consumed",
      "Trigger verification and closure workflow"]),

    ("11-analytics.png", "Repairs — Analytics Dashboard",
     "Data-driven insights into maintenance performance",
     "11",
     ["Visual charts and KPIs for repair operations",
      "Work order completion rates and turnaround times",
      "Technician productivity and resource utilization",
      "Trend analysis for proactive maintenance planning"]),

    ("12-reports.png", "Repairs — Reports",
     "Generate and export comprehensive reports",
     "12",
     ["Pre-built report templates for common needs",
      "Customizable date ranges and filter criteria",
      "Export to PDF and Excel formats",
      "Scheduled report generation and distribution"]),

    ("13-downtime.png", "Repairs — Downtime Tracking",
     "Track and analyze equipment downtime events",
     "13",
     ["Log downtime events linked to maintenance activities",
      "Duration tracking and root cause categorization",
      "Impact analysis on production output",
      "Downtime trends and reduction strategies"]),

    ("14-spare-part-returns.png", "Repairs — Spare Part Returns",
     "Manage return of unused spare parts",
     "14",
     ["Return unused parts from completed work orders",
      "Quality inspection on returned items",
      "Restock to inventory with condition tracking",
      "Cost recovery and waste reduction"]),

    ("15-damaged-tools.png", "Repairs — Damaged Tool Reports",
     "Report and track damaged or worn-out tools",
     "15",
     ["Report tools damaged during repair operations",
      "Damage classification and severity assessment",
      "Replacement request initiation",
      "Vendor warranty claim support documentation"]),
]

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

def add_bg(slide, color=BG_DARK):
    """Add solid background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def add_bottom_bar(slide):
    """Add branded bottom bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.5), SLIDE_W, Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1E)
    shape.line.fill.background()
    add_text_box(slide, Inches(0.5), SLIDE_H - Inches(0.45), Inches(5), Inches(0.4),
                 "iAssetsPro — Enterprise Asset Management Platform",
                 font_size=10, color=DIM, font_name="Calibri")

def add_screenshot(slide, img_path, left, top, width, height):
    """Add screenshot image with rounded-corner card effect."""
    if not os.path.exists(img_path):
        # Fallback: add placeholder
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = RGBColor(0x33, 0x44, 0x55)
        add_text_box(slide, left + Inches(0.2), top + height/2 - Inches(0.2),
                     width - Inches(0.4), Inches(0.4), f"[{img_path} not found]",
                     font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)
        return shape

    # Add image
    pic = slide.shapes.add_picture(img_path, left, top, width, height)

    # Add a subtle border overlay
    border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left - Pt(2), top - Pt(2),
                                     width + Pt(4), height + Pt(4))
    border.fill.background()
    border.line.color.rgb = RGBColor(0x2D, 0x3B, 0x55)
    border.line.width = Pt(2)
    # Move border behind picture
    sp = border._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

    return pic

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Accent line at top
top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
top_line.fill.solid()
top_line.fill.fore_color.rgb = ACCENT
top_line.line.fill.background()

# Title
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1.2),
             "iAssetsPro", font_size=52, color=ACCENT, bold=True, font_name="Calibri Light")

add_text_box(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(1),
             "Maintenance Repairs — Work Order Workflow", font_size=28, color=WHITE, bold=False)

add_text_box(slide, Inches(1.2), Inches(4.0), Inches(10), Inches(0.8),
             "From Operator Request to Completion & Reporting", font_size=18, color=LIGHT_GRAY)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.0), Inches(2), Inches(0.04))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT
div.line.fill.background()

add_text_box(slide, Inches(1.2), Inches(5.3), Inches(10), Inches(0.5),
             "Enterprise Asset Management Platform  |  Lightworld Technology", font_size=14, color=DIM)

# Bottom bar
add_bottom_bar(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: WORKFLOW OVERVIEW (text-based)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Top accent
top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
top_line.fill.solid()
top_line.fill.fore_color.rgb = ACCENT
top_line.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "End-to-End Workflow Overview", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "How a maintenance request flows through the iAssetsPro system from creation to reporting",
             font_size=14, color=DIM)

# Workflow steps
steps = [
    ("1", "Operator\nCreates Request", "Operator identifies issue and submits\nmaintenance request with details", ACCENT),
    ("2", "Supervisor\nReviews & Approves", "Supervisor triages, reviews priority,\nand approves the request", RGBColor(0x06, 0xB6, 0xD4)),
    ("3", "Planner\nCreates Work Order", "Maintenance planner converts approved\nrequest into an executable work order", RGBColor(0x8B, 0x5C, 0xF6)),
    ("4", "Technician\nExecutes WO", "Assigned technician requests tools/materials\nand performs the repair work", RGBColor(0xF5, 0x9E, 0x0B)),
    ("5", "Supervisor\nVerifies & Closes", "Work is verified, completion documented,\nand WO is closed with reporting", RGBColor(0xEF, 0x44, 0x44)),
]

box_w = Inches(2.1)
box_h = Inches(3.6)
start_x = Inches(0.6)
gap = Inches(0.3)
y_pos = Inches(2.0)

for i, (num, title, desc, color) in enumerate(steps):
    x = start_x + i * (box_w + gap)

    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, box_w, box_h)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = RGBColor(0x2D, 0x3B, 0x55)
    card.line.width = Pt(1)

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + box_w/2 - Inches(0.35), y_pos - Inches(0.2),
                                     Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    # Number text
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

    # Connector arrow (except last)
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        x + box_w, y_pos + box_h/2 - Inches(0.1),
                                        gap, Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
        arrow.line.fill.background()

    # Title
    add_text_box(slide, x + Inches(0.15), y_pos + Inches(0.7), box_w - Inches(0.3), Inches(0.8),
                 title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x + Inches(0.15), y_pos + Inches(1.6), box_w - Inches(0.3), Inches(1.8),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ═══════════════════════════════════════════════════════════════════
# CONTENT SLIDES: Each with screenshot + description
# ═══════════════════════════════════════════════════════════════════
for idx, (img_file, title, subtitle, step_num, bullets) in enumerate(slides_data):
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    # Top accent line
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = ACCENT
    top_line.line.fill.background()

    # Step number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.35),
                                    Inches(0.7), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = f"STEP {step_num}"
    p.font.size = Pt(10)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, Inches(1.4), Inches(0.3), Inches(7), Inches(0.5),
                 title, font_size=24, color=WHITE, bold=True)

    # Subtitle
    add_text_box(slide, Inches(1.4), Inches(0.8), Inches(7), Inches(0.4),
                 subtitle, font_size=13, color=ACCENT2)

    # Screenshot on the right (takes ~65% of slide width)
    img_path = os.path.join(SCREENSHOTS, img_file)
    ss_left = Inches(5.0)
    ss_top = Inches(1.4)
    ss_width = Inches(7.8)
    ss_height = Inches(5.4)
    add_screenshot(slide, img_path, ss_left, ss_top, ss_width, ss_height)

    # Bullet points on the left
    bullet_x = Inches(0.6)
    bullet_start_y = Inches(1.6)
    for i, bullet in enumerate(bullets):
        y = bullet_start_y + i * Inches(1.1)
        # Accent dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, bullet_x, y + Pt(5), Pt(8), Pt(8))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        # Text
        add_text_box(slide, bullet_x + Inches(0.3), y, Inches(4.0), Inches(0.9),
                     bullet, font_size=12, color=LIGHT_GRAY)

    add_bottom_bar(slide)

# ═══════════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Top accent
top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
top_line.fill.solid()
top_line.fill.fore_color.rgb = ACCENT
top_line.line.fill.background()

# Centered content
add_text_box(slide, Inches(2), Inches(2.0), Inches(9), Inches(1),
             "Thank You", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

add_text_box(slide, Inches(2), Inches(3.2), Inches(9), Inches(0.6),
             "iAssetsPro — Intelligent Asset Management Platform", font_size=20,
             color=ACCENT, alignment=PP_ALIGN.CENTER)

div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.0), Inches(2), Inches(0.04))
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
div.line.fill.background()

add_text_box(slide, Inches(2), Inches(4.4), Inches(9), Inches(0.5),
             "Lightworld Technology", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
             "https://iassetspro.lightworldtech.com", font_size=13, color=ACCENT2,
             alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ── Save ───────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"✅ Saved: {OUTPUT}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size: {os.path.getsize(OUTPUT) / 1024 / 1024:.1f} MB")