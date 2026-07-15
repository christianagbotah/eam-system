from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import base64, sys

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BG = RGBColor(0x0F, 0x17, 0x2A)
GRN = RGBColor(0x05, 0x96, 0x69)
WHT = RGBColor(0xFF, 0xFF, 0xFF)
GLD = RGBColor(0xF5, 0xA6, 0x23)
LGR = RGBColor(0x94, 0xA3, 0xB8)
DK = RGBColor(0x1E, 0x29, 0x3B)

def bg(s, c=BG):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def rect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background(); return sh

def txt(s, l, t, w, h, text, sz=18, c=WHT, b=False, a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = 'Calibri'; p.alignment = a
    return tb

# SLIDE 1: Title
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s, 0, 0, W, Inches(0.08), GRN)
rect(s, Inches(0.8), Inches(2.5), Inches(0.08), Inches(2.5), GRN)
txt(s, Inches(1.2), Inches(2.3), Inches(10), Inches(0.5), "iAssetsPro", 16, GRN, True)
txt(s, Inches(1.2), Inches(2.7), Inches(10), Inches(0.4), "Enterprise Asset Management Platform", 12, LGR)
txt(s, Inches(1.2), Inches(3.3), Inches(10), Inches(1.5), "Maintenance Repairs\nWork Order Module", 40, WHT, True)
txt(s, Inches(1.2), Inches(5.0), Inches(10), Inches(0.6), "Complete Workflow Presentation — From Request to Closure", 18, GLD)
rect(s, 0, H - Inches(0.5), W, Inches(0.5), RGBColor(0x0A, 0x0F, 0x1F))

# SLIDE 2: Overview
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
rect(s, 0, 0, W, Inches(0.06), GRN)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "WORKFLOW OVERVIEW", 10, GRN, True)
txt(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.6), "End-to-End Maintenance Repair Process", 28, DK, True)
phases = [
    ("Phase 1", "Request & Approval", "• Operator submits\n  Maintenance Request\n• Supervisor reviews\n  & approves\n• Real-time tracking", GRN),
    ("Phase 2", "Work Order Execution", "• Convert MR to WO\n• Assign technician\n• Start work &\n  time tracking", RGBColor(0x0D, 0x94, 0x88)),
    ("Phase 3", "Repair Resources", "• Tool requests\n  & checkout\n• Material requisition\n  & issuance\n• Tool returns", RGBColor(0xF5, 0x9E, 0x0B)),
    ("Phase 4", "Completion & Reporting", "• Work completion\n  sign-off\n• Supervisor closure\n  & lock\n• Analytics & KPIs", RGBColor(0x6D, 0x28, 0xD9)),
]
for i, (ph, ti, desc, co) in enumerate(phases):
    x = Inches(0.6 + i * 3.1)
    rect(s, x, Inches(1.6), Inches(2.8), Inches(4.5), WHT)
    rect(s, x, Inches(1.6), Inches(2.8), Inches(0.08), co)
    txt(s, x+Inches(0.2), Inches(2.0), Inches(2.4), Inches(0.3), ph, 11, co, True)
    txt(s, x+Inches(0.2), Inches(2.4), Inches(2.4), Inches(0.5), ti, 18, DK, True)
    txt(s, x+Inches(0.2), Inches(3.1), Inches(2.4), Inches(3), desc, 12, RGBColor(0x64, 0x74, 0x8B))

# Content slides
slides = [
    (1, "Login & Authentication", "Phase 1 · Request & Approval",
     "Secure login with role-based access control.\n\nSupports multiple user roles:\n• Administrator — Full system access\n• Supervisor — Approve & manage WOs\n• Technician — Execute work orders\n• Operator — Submit maintenance requests\n• Store Keeper — Manage inventory\n• Plant Manager — Plant-wide oversight\n\nEach role has specific permissions for the WO lifecycle."),
    (2, "Dashboard Overview", "Phase 1 · Request & Approval",
     "Comprehensive dashboard after login showing:\n\n• Pending work orders count\n• Open maintenance requests\n• Asset health status indicators\n• Recent activity feed\n• Quick navigation to all modules\n• KPI summary cards\n\nProvides immediate visibility into\nmaintenance operations status."),
    (3, "Maintenance Requests List", "Phase 1 · Request & Approval",
     "All submitted requests in a filterable table:\n\n• Filter by status: Draft, Pending, Approved, Rejected\n• Filter by priority, asset, date range\n• Sort by any column\n• Quick view of request details\n• Bulk actions support\n• Search by request number or description\n\nOperators can create new requests;\nsupervisors can approve/reject from this view."),
    (4, "Maintenance Request Detail", "Phase 1 · Request & Approval",
     "Detailed view of each request:\n\n• Problem description & category\n• Requested priority level\n• Affected asset / equipment\n• Required trades / skills\n• Full approval chain with timestamps\n• Communication thread\n• Attachments & photos\n\nSupervisors can Approve, Reject, or\nRequest More Information."),
    (5, "Work Orders List", "Phase 2 · Work Order Execution",
     "Approved MRs converted to Work Orders:\n\n• Status indicators (Open, In Progress, Completed)\n• Assigned technician\n• Due dates & priority levels\n• Estimated vs actual hours\n• Linked maintenance request\n• Powerful filtering & search\n\nCentral hub for managing all active\nand historical work orders."),
    (6, "Work Order Detail", "Phase 2 · Work Order Execution",
     "Complete WO lifecycle visibility:\n\n• Assigned technician & team\n• Planned vs actual time tracking\n• Task checklist / steps\n• Linked maintenance request\n• Tool & material consumption log\n• Activity history with timestamps\n• Status progression tracker\n• Completion notes & findings"),
    (7, "Tool Requests", "Phase 3 · Repair Resources",
     "Tool management for repair work:\n\n• Request tools needed for WO\n• Status: Pending → Approved → Issued → Returned\n• Full traceability of tool movements\n• Tool shop attendant approval\n• Quantity tracking\n• Linked to specific work orders\n\nEnsures right tools available for\nthe right technician at the right time."),
    (8, "Material Requests", "Phase 3 · Repair Resources",
     "Spare parts & consumables workflow:\n\n• Technician submits material request\n• Supervisor approves\n• Store keeper issues from inventory\n• Consumption tracked against WO\n• Quantity & cost recording\n• Stock level updates\n\nComplete requisition-to-consumption\naudit trail for each part."),
    (9, "Tool Transfers", "Phase 3 · Repair Resources",
     "Tool movement tracking:\n\n• Transfer tools between technicians\n• Return tools to tool shop\n• Chain of custody tracking\n• Who held the tool & when\n• Current holder visibility\n• Damaged tool flagging\n\nMaintains accountability and helps\nplan tool procurement cycles."),
    (10, "Completion & Closure", "Phase 4 · Completion & Reporting",
     "Formal work order completion:\n\n• Technician marks work complete\n• Detailed findings & observations\n• Replaced parts documentation\n• Time log finalization\n• Supervisor/Planner review\n• Formal closure & sign-off\n• WO locked & archived\n\nEnsures quality documentation before\nclosure."),
    (11, "Repair Analytics", "Phase 4 · Completion & Reporting",
     "Real-time operational insights:\n\n• WO completion rate trends\n• Average resolution time\n• Technician productivity metrics\n• Equipment reliability scores\n• MTBF / MTTR analysis\n• Trend analysis for planning\n• Interactive charts & filters\n\nData-driven maintenance decision\nsupport for management."),
    (12, "Repair Reports", "Phase 4 · Completion & Reporting",
     "Comprehensive reporting suite:\n\n• By Asset / Machine reports\n• Overview statistics\n• Technician productivity\n• Materials & cost analysis\n• Downtime impact reports\n• Detailed data exports\n\nExport to Excel, PDF, or Print\nfor stakeholder presentations."),
    (13, "Downtime Tracking", "Phase 4 · Completion & Reporting",
     "Equipment downtime analysis:\n\n• Downtime event recording\n• Impact level classification\n• Duration tracking\n• Affected production lines\n• Root cause analysis\n• Correlation with maintenance\n• Total cost of downtime\n\nCritical for justifying maintenance\ninvestments to management."),
    (14, "Spare Part Returns", "Phase 4 · Completion & Reporting",
     "Unused parts management:\n\n• Return unused parts after WO\n• Inspection & quality check\n• Restock to inventory\n• Credit back to material request\n• Accurate inventory levels\n• Cost reconciliation\n\nPrevents inventory shrinkage and\nensures accurate stock counts."),
    (15, "Damaged Tool Reports", "Phase 4 · Completion & Reporting",
     "Tool damage documentation:\n\n• Report damaged/worn tools\n• Link to specific work orders\n• Damage description & photos\n• Replacement tracking\n• Tool lifecycle data\n• Procurement planning support\n\nHelps maintenance management plan\ntool replacement budgets."),
    (16, "Maintenance Dashboard", "Phase 4 · Completion & Reporting",
     "Command center for maintenance:\n\n• Key performance indicators\n• Pending actions overview\n• Equipment health scores\n• PM compliance rates\n• Maintenance cost trends\n• Team workload distribution\n• Alert & notification center\n\nHolistic view for maintenance\nmanagers and planners."),
]

for num, title, phase, content in slides:
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    rect(s, 0, 0, W, Inches(0.06), GRN)
    rect(s, 0, Inches(0.06), Inches(3.2), H - Inches(0.06), BG)
    txt(s, Inches(0.3), Inches(0.4), Inches(2.6), Inches(0.4), phase.upper(), 10, GRN, True)
    txt(s, Inches(0.3), Inches(0.9), Inches(2.6), Inches(0.8), f"Step {num}", 36, WHT, True)
    txt(s, Inches(0.3), Inches(1.7), Inches(2.6), Inches(1.2), title, 20, WHT, True)
    txt(s, Inches(0.3), Inches(2.9), Inches(2.6), Inches(4), content, 13, LGR)
    # Dots
    for i in range(16):
        x = Inches(0.3 + i * 0.22)
        c = GRN if i == num - 1 else RGBColor(0x33, 0x44, 0x55)
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(6.8), Inches(0.12), Inches(0.12))
        d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background()
    rect(s, 0, H - Inches(0.4), W, Inches(0.4), BG)
    txt(s, Inches(3.5), H - Inches(0.38), Inches(5), Inches(0.3), "iAssetsPro — Maintenance Repairs Work Order Module", 9, LGR)

# Closing
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s, 0, 0, W, Inches(0.08), GRN)
rect(s, Inches(0.8), Inches(2.8), Inches(0.08), Inches(2), GRN)
txt(s, Inches(1.2), Inches(2.6), Inches(10), Inches(0.5), "iAssetsPro", 16, GRN, True)
txt(s, Inches(1.2), Inches(3.3), Inches(10), Inches(1.2), "Thank You", 44, WHT, True)
txt(s, Inches(1.2), Inches(4.6), Inches(10), Inches(0.6), "Ready to transform your maintenance operations", 18, GLD)
txt(s, Inches(1.2), Inches(5.4), Inches(10), Inches(0.8), "Contact: info@iassetspro.com  |  www.lightworldtech.com", 12, LGR)
rect(s, 0, H - Inches(0.5), W, Inches(0.5), RGBColor(0x0A, 0x0F, 0x1F))

out = '/home/z/my-project/public/WO-Workflow-Presentation.pptx'
prs.save(out)
# Base64 encode
with open(out, 'rb') as f:
    b64 = base64.b64encode(f.read()).decode()
print(f"SIZE:{len(b64)}")
print(b64[:100])