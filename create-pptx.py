from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SS_DIR = '/home/z/my-project/download/wo-screenshots'
OUTPUT = '/home/z/my-project/public/WO-Workflow-Presentation.pptx'
os.makedirs('/home/z/my-project/public', exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Dark navy
BG_ACCENT = RGBColor(0x05, 0x96, 0x69)      # Emerald green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
EMERALD = RGBColor(0x05, 0x96, 0x69)
TEAL = RGBColor(0x0D, 0x94, 0x88)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = 0
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_screenshot_slide(prs, title, subtitle, screenshot_file, step_num, total_steps, phase_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_bg(slide, LIGHT_BG)
    
    # Top bar
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.06), EMERALD)
    
    # Left sidebar with step info
    add_rect(slide, 0, Inches(0.06), Inches(3.2), prs.slide_height - Inches(0.06), BG_DARK)
    
    # Phase label
    add_text_box(slide, Inches(0.3), Inches(0.4), Inches(2.6), Inches(0.4),
                 phase_name.upper(), font_size=10, color=EMERALD, bold=True)
    
    # Step number
    add_text_box(slide, Inches(0.3), Inches(0.9), Inches(2.6), Inches(0.8),
                 f"Step {step_num}", font_size=36, color=WHITE, bold=True)
    
    # Title
    add_text_box(slide, Inches(0.3), Inches(1.7), Inches(2.6), Inches(1.2),
                 title, font_size=20, color=WHITE, bold=True)
    
    # Subtitle/description
    add_text_box(slide, Inches(0.3), Inches(2.9), Inches(2.6), Inches(3.5),
                 subtitle, font_size=13, color=LIGHT_GRAY, bold=False)
    
    # Step indicator dots at bottom of sidebar
    dot_y = Inches(6.8)
    dot_start = Inches(0.3)
    for i in range(total_steps):
        x = dot_start + Inches(i * 0.22)
        dot_color = EMERALD if i == step_num - 1 else RGBColor(0x33, 0x44, 0x55)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, dot_y, Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
    
    # Screenshot on the right
    if os.path.exists(screenshot_file):
        slide.shapes.add_picture(screenshot_file, Inches(3.5), Inches(0.5), Inches(9.4), Inches(6.2))
    
    # Bottom bar with branding
    add_rect(slide, 0, prs.slide_height - Inches(0.4), prs.slide_width, Inches(0.4), BG_DARK)
    add_text_box(slide, Inches(3.5), prs.slide_height - Inches(0.38), Inches(5), Inches(0.3),
                 "iAssetsPro — Maintenance Repairs Work Order Module", font_size=9, color=LIGHT_GRAY)
    
    return slide

# =========================================================================
# SLIDE 1: TITLE SLIDE
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

# Accent bar
add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), EMERALD)

# Decorative elements
add_rect(slide, Inches(0.8), Inches(2.5), Inches(0.08), Inches(2.5), EMERALD)

# Logo text
add_text_box(slide, Inches(1.2), Inches(2.3), Inches(10), Inches(0.5),
             "iAssetsPro", font_size=16, color=EMERALD, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.7), Inches(10), Inches(0.4),
             "Enterprise Asset Management Platform", font_size=12, color=LIGHT_GRAY)

# Main title
add_text_box(slide, Inches(1.2), Inches(3.3), Inches(10), Inches(1.5),
             "Maintenance Repairs\nWork Order Module", font_size=40, color=WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(1.2), Inches(5.0), Inches(10), Inches(0.8),
             "Complete Workflow Presentation — From Request to Closure", font_size=18, color=GOLD)

# Bottom info
add_rect(slide, 0, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.5), RGBColor(0x0A, 0x0F, 0x1F))
add_text_box(slide, Inches(1.2), prs.slide_height - Inches(0.45), Inches(5), Inches(0.3),
             "Prepared for Client Presentation", font_size=10, color=LIGHT_GRAY)
add_text_box(slide, Inches(9), prs.slide_height - Inches(0.45), Inches(3.5), Inches(0.3),
             "Lightworld Technologies", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# =========================================================================
# SLIDE 2: WORKFLOW OVERVIEW
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, 0, 0, prs.slide_width, Inches(0.06), EMERALD)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
             "WORKFLOW OVERVIEW", font_size=10, color=EMERALD, bold=True)
add_text_box(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.6),
             "End-to-End Maintenance Repair Process", font_size=28, color=DARK_TEXT, bold=True)

# Four phase boxes
phases = [
    ("Phase 1", "Request & Approval", "• Operator submits Maintenance Request\n• Supervisor reviews & approves\n• Request tracked in real-time", EMERALD),
    ("Phase 2", "Work Order Execution", "• Approved MR converts to Work Order\n• WO assigned to technician\n• Work started with time tracking", TEAL),
    ("Phase 3", "Repair Resources", "• Tool requests & checkout\n• Material requisition & issue\n• Tool returns & transfers", RGBColor(0xF5, 0x9E, 0x0B)),
    ("Phase 4", "Completion & Reporting", "• Work completion sign-off\n• Supervisor closure & lock\n• Analytics & KPI reporting", RGBColor(0x6D, 0x28, 0xD9)),
]

for i, (phase, title, desc, color) in enumerate(phases):
    x = Inches(0.6 + i * 3.1)
    y = Inches(1.6)
    
    # Phase box
    box = add_rect(slide, x, y, Inches(2.8), Inches(4.5), WHITE)
    box.shadow.inherit = False
    
    # Color accent bar at top of box
    add_rect(slide, x, y, Inches(2.8), Inches(0.08), color)
    
    # Phase number
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.3),
                 phase, font_size=11, color=color, bold=True)
    
    # Phase title
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.4), Inches(0.5),
                 title, font_size=18, color=DARK_TEXT, bold=True)
    
    # Phase description
    add_text_box(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.4), Inches(3),
                 desc, font_size=12, color=RGBColor(0x64, 0x74, 0x8B))

# Arrow connectors between phases
for i in range(3):
    x = Inches(3.3 + i * 3.1)
    add_text_box(slide, x, Inches(3.5), Inches(0.4), Inches(0.4),
                 "→", font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide, 0, prs.slide_height - Inches(0.4), prs.slide_width, Inches(0.4), BG_DARK)
add_text_box(slide, Inches(0.6), prs.slide_height - Inches(0.38), Inches(5), Inches(0.3),
             "iAssetsPro — Maintenance Repairs Work Order Module", font_size=9, color=LIGHT_GRAY)

# =========================================================================
# SCREENSHOT SLIDES
# =========================================================================
slides_data = [
    # Phase 1: Request & Approval
    ("Login & Authentication", 
     "Users authenticate via a secure login page with role-based access control. The system supports multiple user roles including Administrators, Supervisors, Technicians, Operators, and Store Keepers — each with specific permissions for the WO lifecycle.",
     "01-login-page.png", 1, 16, "Phase 1 · Request & Approval"),
    ("Dashboard Overview", 
     "After login, users are presented with a comprehensive dashboard showing key metrics: pending work orders, open maintenance requests, asset health status, and recent activity. The dashboard provides quick navigation to all modules.",
     "02-dashboard.png", 2, 16, "Phase 1 · Request & Approval"),
    ("Maintenance Requests List", 
     "The Maintenance Requests module displays all submitted requests in a filterable, sortable table. Users can filter by status (Draft, Pending, Approved, Rejected), priority, asset, and date range. The table shows request details at a glance.",
     "03-maintenance-requests-list.png", 3, 16, "Phase 1 · Request & Approval"),
    ("Maintenance Request Detail", 
     "Each Maintenance Request has a detailed view showing the problem description, requested priority, affected asset/equipment, required trades, and the full approval chain with timestamps. Authorized users can approve, reject, or request more information.",
     "05-mr-detail.png", 4, 16, "Phase 1 · Request & Approval"),
    
    # Phase 2: Work Order Execution
    ("Work Orders List", 
     "Approved Maintenance Requests are converted into Work Orders. The WO list provides a comprehensive view of all work orders with status indicators, assigned technicians, due dates, and priority levels. Powerful filtering and search capabilities help users find specific orders.",
     "06-work-orders-list.png", 5, 16, "Phase 2 · Work Order Execution"),
    ("Work Order Detail", 
     "The Work Order detail page provides complete visibility into the work order lifecycle — including assigned technician, planned vs actual hours, task checklist, linked maintenance request, tool/material consumption, and all activity history with timestamps.",
     "07-wo-detail.png", 6, 16, "Phase 2 · Work Order Execution"),
    
    # Phase 3: Repair Resources
    ("Tool Requests", 
     "Technicians can request tools needed for their assigned work orders. The Tool Requests module tracks each request's status — from submission through approval, issuance by the tool shop, and return. Full traceability of tool movements is maintained.",
     "08-tool-requests.png", 7, 16, "Phase 3 · Repair Resources"),
    ("Material Requests", 
     "The Material Requests module handles spare parts and consumable requisitions for repair work. Requests flow from technician submission → supervisor approval → store keeper issuance → consumption tracking against the work order.",
     "09-material-requests.png", 8, 16, "Phase 3 · Repair Resources"),
    ("Tool Transfers", 
     "Tools can be transferred between technicians or returned to the tool shop. The Tool Transfers module tracks the complete chain of custody — who held the tool, when it was transferred, and the current holder. Damaged tools are flagged for reporting.",
     "10-tool-transfers.png", 9, 16, "Phase 3 · Repair Resources"),
    
    # Phase 4: Completion & Reporting
    ("Completion & Closure", 
     "The Completion & Closure module allows technicians to mark work as complete with detailed findings, replaced parts, and time logs. Supervisors and planners review and formally close the work order, which then becomes locked and archived.",
     "11-completion-closure.png", 10, 16, "Phase 4 · Completion & Reporting"),
    ("Repair Analytics", 
     "The Analytics dashboard provides real-time insights into repair operations — including WO completion rates, average resolution times, technician productivity, equipment reliability metrics, and trend analysis for proactive maintenance planning.",
     "12-repair-analytics.png", 11, 16, "Phase 4 · Completion & Reporting"),
    ("Repair Reports", 
     "Comprehensive reporting module with multiple report types: by asset/machine, overview statistics, technician productivity, materials & cost analysis, downtime impact, and detailed data exports to Excel, PDF, and print formats.",
     "13-repair-reports.png", 12, 16, "Phase 4 · Completion & Reporting"),
    ("Downtime Tracking", 
     "Track and analyze equipment downtime events with impact level classification. The module records downtime duration, affected production lines, root causes, and correlates with maintenance activities for total cost of downtime analysis.",
     "14-downtime-tracking.png", 13, 16, "Phase 4 · Completion & Reporting"),
    ("Spare Part Returns", 
     "Manage the return of unused spare parts and consumables after work order completion. Parts are inspected, restocked to inventory, and credited back to the original material request, maintaining accurate inventory levels.",
     "15-spare-part-returns.png", 14, 16, "Phase 4 · Completion & Reporting"),
    ("Damaged Tool Reports", 
     "Document and track damaged or worn tools discovered during repair work. The system creates records for tool replacement, links damage to specific work orders, and helps maintenance management plan tool procurement and replacement cycles.",
     "16-damaged-tool-reports.png", 15, 16, "Phase 4 · Completion & Reporting"),
    ("Maintenance Dashboard", 
     "The dedicated Maintenance Dashboard provides a holistic view of all maintenance operations — KPIs, pending actions, equipment health scores, PM compliance rates, and maintenance cost trends. It serves as the command center for maintenance management.",
     "17-maintenance-dashboard.png", 16, 16, "Phase 4 · Completion & Reporting"),
]

for title, desc, ss_file, step, total, phase in slides_data:
    add_screenshot_slide(prs, title, desc, os.path.join(SS_DIR, ss_file), step, total, phase)

# =========================================================================
# CLOSING SLIDE
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), EMERALD)

# Decorative line
add_rect(slide, Inches(0.8), Inches(2.8), Inches(0.08), Inches(2), EMERALD)

add_text_box(slide, Inches(1.2), Inches(2.6), Inches(10), Inches(0.5),
             "iAssetsPro", font_size=16, color=EMERALD, bold=True)

add_text_box(slide, Inches(1.2), Inches(3.3), Inches(10), Inches(1.2),
             "Thank You", font_size=44, color=WHITE, bold=True)

add_text_box(slide, Inches(1.2), Inches(4.6), Inches(10), Inches(0.6),
             "Ready to transform your maintenance operations", font_size=18, color=GOLD)

add_text_box(slide, Inches(1.2), Inches(5.4), Inches(10), Inches(0.8),
             "Contact: info@iassetspro.com  |  www.lightworldtech.com", font_size=12, color=LIGHT_GRAY)

add_rect(slide, 0, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.5), RGBColor(0x0A, 0x0F, 0x1F))
add_text_box(slide, Inches(1.2), prs.slide_height - Inches(0.45), Inches(5), Inches(0.3),
             "© 2025 Lightworld Technologies. All rights reserved.", font_size=10, color=LIGHT_GRAY)

# =========================================================================
# SAVE
# =========================================================================
prs.save(OUTPUT)
print(f"✅ Presentation saved to: {OUTPUT}")
print(f"   Total slides: {len(prs.slides)}")