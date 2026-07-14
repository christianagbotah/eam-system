from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

SS = '/home/z/my-project/download/wo-screenshots'
OUT = '/home/z/my-project/public/WO-Workflow-Presentation.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BG = RGBColor(0x0F,0x17,0x2A); GRN = RGBColor(0x05,0x96,0x69); WHT = RGBColor(0xFF,0xFF,0xFF)
GLD = RGBColor(0xF5,0xA6,0x23); LGR = RGBColor(0x94,0xA3,0xB8); DK = RGBColor(0x1E,0x29,0x3B)
LBG = RGBColor(0xF8,0xFA,0xFC)

def bg(s,c=BG):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def rect(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background(); return sh

def txt(s,l,t,w,h,text,sz=18,c=WHT,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name='Calibri'; p.alignment=a
    return tb

def ss_slide(prs, title, desc, img, step, total, phase):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    rect(s,0,0,W,Inches(0.06),GRN)
    rect(s,0,Inches(0.06),Inches(3.2),H-Inches(0.06),BG)
    txt(s,Inches(0.3),Inches(0.4),Inches(2.6),Inches(0.4),phase.upper(),10,GRN,True)
    txt(s,Inches(0.3),Inches(0.9),Inches(2.6),Inches(0.8),f"Step {step}",36,WHT,True)
    txt(s,Inches(0.3),Inches(1.7),Inches(2.6),Inches(1.2),title,20,WHT,True)
    # Description with word wrap
    tb = s.shapes.add_textbox(Inches(0.3), Inches(3.0), Inches(2.6), Inches(3.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(desc.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line; p.font.size = Pt(11); p.font.color.rgb = LGR; p.font.name = 'Calibri'
        p.space_after = Pt(2)
    # Progress dots
    for i in range(total):
        x = Inches(0.3 + i * 0.18)
        dc = GRN if i == step-1 else RGBColor(0x33,0x44,0x55)
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(6.9), Inches(0.1), Inches(0.1))
        d.fill.solid(); d.fill.fore_color.rgb = dc; d.line.fill.background()
    # Screenshot
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(3.5), Inches(0.3), Inches(9.5), H-Inches(0.8))
    # Footer
    rect(s,0,H-Inches(0.35),W,Inches(0.35),BG)
    txt(s,Inches(3.5),H-Inches(0.32),Inches(5),Inches(0.3),"iAssetsPro — Maintenance Repairs Work Order Module",8,LGR)
    return s

# ===== SLIDE 1: TITLE =====
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,W,Inches(0.08),GRN)
rect(s,Inches(0.8),Inches(2.5),Inches(0.08),Inches(2.5),GRN)
txt(s,Inches(1.2),Inches(2.3),Inches(10),Inches(0.5),"iAssetsPro",16,GRN,True)
txt(s,Inches(1.2),Inches(2.7),Inches(10),Inches(0.4),"Enterprise Asset Management Platform",12,LGR)
txt(s,Inches(1.2),Inches(3.3),Inches(10),Inches(1.5),"Maintenance Repairs\nWork Order Module",40,WHT,True)
txt(s,Inches(1.2),Inches(5.0),Inches(10),Inches(0.6),"Complete Workflow Presentation — From Request to Closure",18,GLD)
rect(s,0,H-Inches(0.5),W,Inches(0.5),RGBColor(0x0A,0x0F,0x1F))
txt(s,Inches(1.2),H-Inches(0.45),Inches(5),Inches(0.3),"Prepared for Client Presentation",10,LGR)
txt(s,Inches(9),H-Inches(0.45),Inches(3.5),Inches(0.3),"Lightworld Technologies",10,LGR, False, PP_ALIGN.RIGHT)

# ===== SLIDE 2: OVERVIEW =====
s=prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = LBG
rect(s,0,0,W,Inches(0.06),GRN)
txt(s,Inches(0.6),Inches(0.3),Inches(12),Inches(0.5),"WORKFLOW OVERVIEW",10,GRN,True)
txt(s,Inches(0.6),Inches(0.7),Inches(12),Inches(0.6),"End-to-End Maintenance Repair Process",28,DK,True)
phases=[("Phase 1","Request & Approval","• Operator submits\n  Maintenance Request\n• Supervisor reviews\n  & approves request\n• Real-time status tracking",GRN),
("Phase 2","Work Order Execution","• Convert approved MR\n  to Work Order\n• Assign to technician\n• Start work & log time\n• Track progress",RGBColor(0x0D,0x94,0x88)),
("Phase 3","Repair Resources","• Tool request & checkout\n• Material requisition\n• Tool transfers\n• Downtime logging\n• Inventory tracking",RGBColor(0xF5,0x9E,0x0B)),
("Phase 4","Completion & Reporting","• Work completion sign-off\n• Supervisor closure\n• WO locked & archived\n• Analytics & KPIs\n• Compliance reports",RGBColor(0x6D,0x28,0xD9))]
for i,(ph,ti,desc,co) in enumerate(phases):
    x=Inches(0.6+i*3.1)
    rect(s,x,Inches(1.6),Inches(2.8),Inches(4.8),WHT)
    rect(s,x,Inches(1.6),Inches(2.8),Inches(0.08),co)
    txt(s,x+Inches(0.2),Inches(2.0),Inches(2.4),Inches(0.3),ph,11,co,True)
    txt(s,x+Inches(0.2),Inches(2.4),Inches(2.4),Inches(0.5),ti,18,DK,True)
    tb=s.shapes.add_textbox(x+Inches(0.2),Inches(3.1),Inches(2.4),Inches(3))
    tf=tb.text_frame; tf.word_wrap=True
    for j,line in enumerate(desc.split('\n')):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
        p.text=line; p.font.size=Pt(12); p.font.color.rgb=RGBColor(0x64,0x74,0x8B); p.font.name='Calibri'
    if i<3:
        txt(s,Inches(3.3+i*3.1),Inches(3.6),Inches(0.4),Inches(0.4),"→",24,LGR,False,PP_ALIGN.CENTER)
rect(s,0,H-Inches(0.4),W,Inches(0.4),BG)
txt(s,Inches(0.6),H-Inches(0.38),Inches(5),Inches(0.3),"iAssetsPro — Maintenance Repairs Work Order Module",9,LGR)

# ===== SCREENSHOT SLIDES =====
slides = [
    ("Dashboard","Central command center showing pending WOs, maintenance requests, asset health, and quick navigation to all modules.","02-dashboard.png",1,17,"Phase 1 · Request & Approval"),
    ("Maintenance Requests","All submitted requests in a filterable table. Filter by status, priority, asset, and date. Supervisors approve/reject from here.","02-maintenance-requests.png",2,17,"Phase 1 · Request & Approval"),
    ("Create New Request","Operators fill out the request form: describe the problem, select asset, set priority, and choose required trades.","03-create-mr-dialog.png",3,17,"Phase 1 · Request & Approval"),
    ("Request Detail","Full view of each request showing description, priority, asset details, approval chain with timestamps, and communication thread.","04-mr-detail.png",4,17,"Phase 1 · Request & Approval"),
    ("Work Orders List","Approved MRs converted to Work Orders. Track status, assigned technician, due dates, and priority levels at a glance.","05-work-orders.png",5,17,"Phase 2 · Execution"),
    ("Work Order Detail","Complete WO lifecycle: assigned technician, task progress, time tracking, linked MR, and action buttons for tools, materials, and completion.","06-wo-detail.png",6,17,"Phase 2 · Execution"),
    ("Tool Request","From the WO detail, technicians request tools needed for the repair. Requests go to the tool shop for approval and issuance.","07-tool-request-dialog.png",7,17,"Phase 3 · Repair Resources"),
    ("Material Request","Request spare parts and consumables. The requisition flows from technician → supervisor → store keeper → issuance.","08-material-request-dialog.png",8,17,"Phase 3 · Repair Resources"),
    ("Tool Transfer","Transfer tools between technicians or return them to the tool shop. Full chain-of-custody tracking maintained.","09-tool-transfer-dialog.png",9,17,"Phase 3 · Repair Resources"),
    ("Complete Work Order","Mark work complete with findings, replaced parts, and final time log. Triggers supervisor review workflow.","10-complete-wo-dialog.png",10,17,"Phase 3 · Repair Resources"),
    ("Tool Requests List","Central view of all tool requests across work orders. Track status from submission through issuance to return.","11-tool-requests.png",11,17,"Phase 3 · Repair Resources"),
    ("Material Requests List","All material requisitions in one place. Monitor approval status, issuance progress, and consumption tracking.","12-material-requests.png",12,17,"Phase 4 · Completion & Reporting"),
    ("Tool Transfers","Track all tool movements between technicians and the tool shop. Ensures accountability and proper tool management.","13-tool-transfers.png",13,17,"Phase 4 · Completion & Reporting"),
    ("Downtime Tracking","Record and analyze equipment downtime. Classify impact levels and correlate with maintenance activities for cost analysis.","14-downtime.png",14,17,"Phase 4 · Completion & Reporting"),
    ("Completion & Closure","Formal WO completion review. Supervisors verify work quality, close the order, and lock it for archival.","15-completion.png",15,17,"Phase 4 · Completion & Reporting"),
    ("Spare Part Returns","Return unused parts after WO completion. Parts are inspected, restocked, and credited back to the original request.","16-spare-part-returns.png",16,17,"Phase 4 · Completion & Reporting"),
    ("Damaged Tool Reports","Document and track damaged tools discovered during repair. Supports replacement planning and procurement budgets.","17-damaged-tools.png",17,17,"Phase 4 · Completion & Reporting"),
]

for title,desc,img,step,total,phase in slides:
    ss_slide(prs, title, desc, os.path.join(SS, img), step, total, phase)

# Analytics slide
ss_slide(prs, "Repair Analytics","Real-time insights: WO completion rates, resolution times, technician productivity, equipment reliability, MTBF/MTTR, and trend analysis.",
    os.path.join(SS, "18-analytics.png"), 18, 18, "Phase 4 · Completion & Reporting")

# Reports slide
ss_slide(prs, "Repair Reports","Comprehensive reporting: by asset, overview stats, technician productivity, materials & costs, downtime impact. Export to Excel, PDF, or Print.",
    os.path.join(SS, "19-repair-reports.png"), 19, 19, "Phase 4 · Completion & Reporting")

# ===== CLOSING =====
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,W,Inches(0.08),GRN)
rect(s,Inches(0.8),Inches(2.8),Inches(0.08),Inches(2),GRN)
txt(s,Inches(1.2),Inches(2.6),Inches(10),Inches(0.5),"iAssetsPro",16,GRN,True)
txt(s,Inches(1.2),Inches(3.3),Inches(10),Inches(1.2),"Thank You",44,WHT,True)
txt(s,Inches(1.2),Inches(4.6),Inches(10),Inches(0.6),"Ready to transform your maintenance operations",18,GLD)
txt(s,Inches(1.2),Inches(5.4),Inches(10),Inches(0.8),"Contact: info@iassetspro.com  |  www.lightworldtech.com",12,LGR)
rect(s,0,H-Inches(0.5),W,Inches(0.5),RGBColor(0x0A,0x0F,0x1F))
txt(s,Inches(1.2),H-Inches(0.45),Inches(5),Inches(0.3),"© 2025 Lightworld Technologies. All rights reserved.",10,LGR)

prs.save(OUT)
print(f"✅ Saved: {OUT} ({os.path.getsize(OUT)//1024}KB, {len(prs.slides)} slides)")